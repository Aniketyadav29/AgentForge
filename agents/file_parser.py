"""
AgentForge — Universal File Parser
Handles PDF, CSV, Excel, DOCX, TXT/Markdown/JSON/HTML, and common image files.
Returns structured text + DataFrames for downstream processing.
"""

import io
import os
import json
import re
from pathlib import Path
from typing import Optional
import pandas as pd


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def _safe_import(module: str):
    """Try to import a module and return None on failure."""
    import importlib
    try:
        return importlib.import_module(module)
    except ImportError:
        return None


# ─────────────────────────────────────────────────────────────────────────────
# Individual parsers
# ─────────────────────────────────────────────────────────────────────────────

def _parse_pdf(file_bytes: bytes) -> dict:
    """Parse PDF using pdfplumber for accurate text + table extraction."""
    pdfplumber = _safe_import("pdfplumber")
    if not pdfplumber:
        raise ImportError("pdfplumber not installed. Run: pip install pdfplumber")

    text_pages = []
    tables: list[pd.DataFrame] = []

    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            # Extract plain text
            page_text = page.extract_text() or ""
            text_pages.append(f"[Page {page_num}]\n{page_text}")

            # Extract tables as DataFrames
            for table in page.extract_tables():
                if table and len(table) > 1:
                    try:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        df = df.dropna(how="all")
                        tables.append(df)
                    except Exception:
                        pass

    return {
        "text": "\n\n".join(text_pages),
        "tables": tables,
        "page_count": len(pdf.pages) if hasattr(pdf, "pages") else None,
    }


def _parse_csv(file_bytes: bytes) -> dict:
    """Parse CSV using pandas."""
    try:
        df = pd.read_csv(io.BytesIO(file_bytes))
    except Exception:
        df = pd.read_csv(io.BytesIO(file_bytes), encoding="latin-1")

    text = f"CSV Data ({len(df)} rows × {len(df.columns)} columns)\n\n"
    text += f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
    text += "Data Preview (first 20 rows):\n"
    text += df.head(20).to_string(index=False)
    text += "\n\nStatistical Summary:\n"
    text += df.describe(include="all").to_string()

    return {"text": text, "tables": [df], "row_count": len(df)}


def _parse_excel(file_bytes: bytes) -> dict:
    """Parse Excel (.xlsx / .xls) — all sheets."""
    xl = pd.ExcelFile(io.BytesIO(file_bytes))
    all_text = []
    tables = []

    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        tables.append(df)
        text = f"Sheet: {sheet_name} ({len(df)} rows × {len(df.columns)} cols)\n"
        text += f"Columns: {', '.join(str(c) for c in df.columns)}\n"
        text += df.head(20).to_string(index=False)
        all_text.append(text)

    return {
        "text": "\n\n".join(all_text),
        "tables": tables,
        "sheet_count": len(xl.sheet_names),
    }


def _parse_docx(file_bytes: bytes) -> dict:
    """Parse DOCX using python-docx."""
    docx = _safe_import("docx")
    if not docx:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

    from docx import Document
    doc = Document(io.BytesIO(file_bytes))

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for tbl in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in tbl.rows]
        if rows:
            try:
                df = pd.DataFrame(rows[1:], columns=rows[0])
                tables.append(df)
            except Exception:
                pass

    return {
        "text": "\n".join(paragraphs),
        "tables": tables,
        "paragraph_count": len(paragraphs),
    }


def _parse_txt(file_bytes: bytes) -> dict:
    """Parse plain text files."""
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1")
    return {"text": text, "tables": []}


def _parse_json(file_bytes: bytes) -> dict:
    """Parse JSON as readable structured text."""
    try:
        raw = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raw = file_bytes.decode("latin-1")
    try:
        data = json.loads(raw)
        text = json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        text = raw
    return {"text": f"JSON document\n\n{text}", "tables": []}


def _parse_html(file_bytes: bytes) -> dict:
    """Parse HTML with BeautifulSoup when available, otherwise strip tags."""
    try:
        raw = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raw = file_bytes.decode("latin-1")

    bs4 = _safe_import("bs4")
    if bs4:
        soup = bs4.BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n")
    else:
        text = re.sub(r"<[^>]+>", " ", raw)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return {"text": text.strip(), "tables": []}


def _parse_pptx(file_bytes: bytes) -> dict:
    """Parse PowerPoint text when python-pptx is installed."""
    pptx = _safe_import("pptx")
    if not pptx:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    from pptx import Presentation
    deck = Presentation(io.BytesIO(file_bytes))
    slides = []
    for index, slide in enumerate(deck.slides, 1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        slides.append(f"[Slide {index}]\n" + "\n".join(lines))
    return {"text": "\n\n".join(slides), "tables": [], "slide_count": len(deck.slides)}


def _parse_image(file_bytes: bytes) -> dict:
    """Parse image metadata and OCR text if optional local OCR tooling is available."""
    metadata = []
    ocr_text = ""

    image = None
    pil = _safe_import("PIL.Image")
    if pil:
        try:
            image = pil.open(io.BytesIO(file_bytes))
            metadata.append(f"Dimensions: {image.width} x {image.height}px")
            metadata.append(f"Color mode: {image.mode}")
            metadata.append(f"Format: {image.format or 'unknown'}")
        except Exception as exc:
            metadata.append(f"Image metadata extraction failed: {exc}")

    pytesseract = _safe_import("pytesseract")
    if pytesseract and image is not None:
        try:
            ocr_text = pytesseract.image_to_string(image).strip()
        except Exception as exc:
            metadata.append(f"OCR unavailable: {exc}")
    else:
        metadata.append("OCR unavailable: install Pillow + pytesseract + Tesseract OCR for text extraction.")

    text_parts = [
        "Image analysis document",
        "\n".join(metadata) if metadata else "No image metadata available.",
    ]
    if ocr_text:
        text_parts.append(f"Detected image text:\n{ocr_text}")
    else:
        text_parts.append("Detected image text: none extracted.")

    return {"text": "\n\n".join(text_parts), "tables": [], "is_image": True}


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".pdf", ".csv", ".xlsx", ".xls", ".docx", ".txt", ".md",
    ".json", ".html", ".htm", ".pptx",
    ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif",
}


def parse_file(file_bytes: bytes, filename: str) -> dict:
    """
    Parse an uploaded file and return structured data.

    Args:
        file_bytes: Raw bytes of the uploaded file.
        filename:   Original filename (used to detect format).

    Returns:
        {
            "text":      Full extracted text (str),
            "tables":    List of pandas DataFrames extracted from the file,
            "metadata":  Dict with file-specific metadata,
            "file_type": Detected file type string,
        }

    Raises:
        ValueError if the file extension is not supported.
    """
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    parsers = {
        ".pdf":  _parse_pdf,
        ".csv":  _parse_csv,
        ".xlsx": _parse_excel,
        ".xls":  _parse_excel,
        ".docx": _parse_docx,
        ".txt":  _parse_txt,
        ".md":   _parse_txt,
        ".json": _parse_json,
        ".html": _parse_html,
        ".htm":  _parse_html,
        ".pptx": _parse_pptx,
        ".png":  _parse_image,
        ".jpg":  _parse_image,
        ".jpeg": _parse_image,
        ".webp": _parse_image,
        ".bmp":  _parse_image,
        ".gif":  _parse_image,
    }

    result = parsers[ext](file_bytes)
    result["file_type"] = ext.lstrip(".")
    result["metadata"] = {
        "filename": filename,
        "file_type": ext.lstrip("."),
        "file_size_bytes": len(file_bytes),
        "table_count": len(result.get("tables", [])),
        "has_tables": len(result.get("tables", [])) > 0,
    }
    return result
